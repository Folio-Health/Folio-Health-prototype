"""Regenerate EMR-Working-Logic-Research.pptx from the three research reports.

Run from the folio-health directory:  python docs/research/build-deck.py
Requires python-pptx (`pip install python-pptx`).

The palette is the app's own design system (src/app/globals.css):
Folio Deep #0E443B · Folio Teal #11534A · Signal Coral #E65A4F (small accents
only — brand guideline: "accent only, never large areas") · Mint #7FC9BC for
secondary text on dark brand surfaces · pale teal-gray surfaces · Poppins type.
"""

import re
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

DOCS = [
    ("docs/research/emr-front-office.md", "Front-Office Workflows",
     "Appointments · Registration & Identity · Check-In & Queuing · Walk-Ins & Triage"),
    ("docs/research/emr-clinical-flows.md", "Clinical Workflows",
     "Encounters & Notes · CPOE · Closed-Loop Medication · Nursing · ADT · Referrals"),
    ("docs/research/emr-back-office-security.md", "Back-Office & Security",
     "Revenue Cycle · HIM · Security & Break-Glass · Account Lifecycle · Inventory · Messaging"),
]

# ── Folio design tokens (src/app/globals.css) ────────────────────────────────
DEEP = RGBColor(0x0E, 0x44, 0x3B)        # --primary / --foreground  (Folio Deep)
DEEP_LOW = RGBColor(0x0C, 0x3E, 0x36)    # --stats-bg (Folio Deep low)
DEEP_NIGHT = RGBColor(0x07, 0x1F, 0x1A)  # --hero-gradient-start
TEAL = RGBColor(0x11, 0x53, 0x4A)        # --chart-1 (Folio Teal)
CORAL = RGBColor(0xE6, 0x5A, 0x4F)       # --brand-coral (Signal Coral)
MINT = RGBColor(0x7F, 0xC9, 0xBC)        # --brand-mint
SURFACE = RGBColor(0xFF, 0xFF, 0xFF)     # --background
MUTED_BG = RGBColor(0xF3, 0xF6, 0xF5)    # --muted-bg
BORDER = RGBColor(0xDC, 0xE3, 0xE1)      # --border
MUTED_FG = RGBColor(0x5C, 0x6D, 0x69)    # --muted-foreground

SANS = "Poppins"
MONO = "Consolas"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(BLANK)


def fill_rect(slide, x, y, w, h, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def set_bg(slide, color):
    fill_rect(slide, 0, 0, prs.slide_width, prs.slide_height, color)


def para(tf, first):
    return tf.paragraphs[0] if first else tf.add_paragraph()


def run(p, text, size, color, font=SANS, bold=False):
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = font
    return r


def footer(slide, dark=False):
    # coral dot + wordmark, echoing the app's sidebar brand row
    dot = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(0.55), Inches(7.13), Inches(0.09), Inches(0.09))
    dot.fill.solid()
    dot.fill.fore_color.rgb = CORAL
    dot.line.fill.background()
    dot.shadow.inherit = False
    tb = slide.shapes.add_textbox(Inches(0.72), Inches(6.98), Inches(8), Inches(0.4))
    p = tb.text_frame.paragraphs[0]
    run(p, "Folio Health · EMR Working-Logic Research",
        9.5, MINT if dark else MUTED_FG)


def title_bar(slide, text, sub=None):
    tb = slide.shapes.add_textbox(Inches(0.55), Inches(0.3), Inches(12.2), Inches(0.95))
    tf = tb.text_frame
    tf.word_wrap = True
    run(para(tf, True), text, 23, DEEP, bold=True)
    if sub:
        run(tf.add_paragraph(), sub, 11.5, MUTED_FG)
    # thin Folio Teal rule under the header, coral tick at its start
    fill_rect(slide, Inches(0.58), Inches(1.24), Inches(0.45), Pt(2.6), CORAL)
    fill_rect(slide, Inches(1.03), Inches(1.24), Inches(11.7), Pt(2.6), BORDER)


def strip_md(text):
    text = re.sub(r"\[([^\]]+)\]\(([^)]+)\)", r"\1", text)
    return text.replace("**", "").replace("*", "").replace("`", "")


LINES_PER_SLIDE = 20


def render_section(doc_title, sec_title, body_lines):
    pages, cur, count = [], [], 0
    for ln in body_lines:
        weight = 1 + (len(strip_md(ln)) // 110)
        if count + weight > LINES_PER_SLIDE and cur:
            pages.append(cur)
            cur, count = [], 0
        cur.append(ln)
        count += weight
    if cur:
        pages.append(cur)

    for i, page in enumerate(pages):
        slide = add_slide()
        set_bg(slide, SURFACE)
        t = sec_title if len(pages) == 1 else f"{sec_title}  ({i + 1}/{len(pages)})"
        title_bar(slide, t, doc_title)
        footer(slide)

        # code-block panel styling: put code lines on a muted panel
        tb = slide.shapes.add_textbox(Inches(0.55), Inches(1.42), Inches(12.25), Inches(5.45))
        tf = tb.text_frame
        tf.word_wrap = True
        first, in_code = True, False
        for ln in page:
            s = ln.rstrip()
            if s.strip().startswith("```"):
                in_code = not in_code
                continue
            p = para(tf, first)
            first = False
            if in_code:
                r = run(p, ("   " + s) if s.strip() else " ", 11.5, TEAL, font=MONO)
                continue
            if s.startswith("#### ") or s.startswith("### "):
                p.space_before = Pt(8)
                run(p, s.lstrip("# ").strip(), 14.5, TEAL, bold=True)
                continue
            if s.strip().startswith("|"):
                run(p, strip_md(s), 10.5, DEEP, font=MONO)
                continue
            m = re.match(r"^(\s*)([-*]|\d+\.)\s+(.*)$", s)
            if m:
                indent, marker, rest = m.groups()
                p.level = min(len(indent) // 2, 3)
                bullet = marker + " " if marker.rstrip(".").isdigit() else "•  "
                run(p, bullet + strip_md(rest), 12.5, DEEP)
                continue
            run(p, strip_md(s) if s.strip() else " ", 12.5, DEEP)


def divider(title, subtitle):
    slide = add_slide()
    set_bg(slide, DEEP_LOW)
    # deep night band at the top, like the hero gradient
    fill_rect(slide, 0, 0, prs.slide_width, Inches(2.2), DEEP_NIGHT)
    fill_rect(slide, Inches(0.92), Inches(3.6), Inches(0.7), Pt(3), CORAL)
    tb = slide.shapes.add_textbox(Inches(0.9), Inches(2.75), Inches(11.5), Inches(2.4))
    tf = tb.text_frame
    tf.word_wrap = True
    run(para(tf, True), title, 36, SURFACE, bold=True)
    p = tf.add_paragraph()
    p.space_before = Pt(14)
    run(p, subtitle, 14.5, MINT)
    footer(slide, dark=True)


# ── cover ────────────────────────────────────────────────────────────────────
slide = add_slide()
set_bg(slide, DEEP_NIGHT)
fill_rect(slide, 0, Inches(5.9), prs.slide_width, Inches(1.6), DEEP_LOW)
fill_rect(slide, Inches(0.92), Inches(1.95), Inches(0.7), Pt(3), CORAL)
tb = slide.shapes.add_textbox(Inches(0.9), Inches(2.15), Inches(11.6), Inches(3.4))
tf = tb.text_frame
tf.word_wrap = True
run(para(tf, True), "EMR Working-Logic Research", 42, SURFACE, bold=True)
p = tf.add_paragraph()
p.space_before = Pt(12)
run(p, "How production EMRs — Epic, Oracle Cerner, athenahealth, MEDITECH, "
       "OpenMRS, Medplum — design every day-to-day hospital workflow, "
       "appointments to sign-out.", 16, MINT)
p = tf.add_paragraph()
p.space_before = Pt(10)
run(p, "Three reports · state machines, actors, business rules, FHIR mappings "
       "· sources cited in docs/research/*.md", 12.5, RGBColor(0xB8, 0xD8, 0xD1))
tb = slide.shapes.add_textbox(Inches(0.9), Inches(6.3), Inches(11.5), Inches(0.8))
run(tb.text_frame.paragraphs[0], "Folio Health · August 2026", 12, MINT)

# ── reports ──────────────────────────────────────────────────────────────────
for path, doc_title, subtitle in DOCS:
    text = open(path, encoding="utf-8").read()
    divider(doc_title, subtitle)

    sections, cur_title, cur_body = [], None, []
    for ln in text.split("\n"):
        if ln.startswith("## "):
            if cur_title:
                sections.append((cur_title, cur_body))
            cur_title, cur_body = strip_md(ln[3:]).strip(), []
        elif cur_title is not None:
            cur_body.append(ln)
    if cur_title:
        sections.append((cur_title, cur_body))

    for sec_title, body in sections:
        out, prev_blank, skip = [], False, False
        for ln in body:
            if ln.startswith("### Sources"):
                skip = True
                continue
            if ln.startswith("### "):
                skip = False
            if skip:
                continue
            blank = not ln.strip()
            if blank and prev_blank:
                continue
            prev_blank = blank
            out.append(ln)
        while out and not out[0].strip():
            out.pop(0)
        while out and not out[-1].strip():
            out.pop()
        if out:
            render_section(doc_title, sec_title, out)

prs.save("docs/research/EMR-Working-Logic-Research.pptx")
print("slides:", len(prs.slides._sldIdLst))
